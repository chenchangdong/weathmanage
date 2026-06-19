"""资产配置报告 PPT 导出（封面 + 目录 + 章节空白页）。"""

from __future__ import annotations

import io
from datetime import date
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 报告章节（与模板目录页一致）
REPORT_CHAPTERS: list[dict[str, str]] = [
    {"id": "01", "title": "总览"},
    {"id": "02", "title": "资产构成"},
    {"id": "03", "title": "持仓穿透"},
    {"id": "04", "title": "收益归因"},
    {"id": "05", "title": "风险控制"},
    {"id": "06", "title": "客户行为"},
    {"id": "07", "title": "相关词语释义"},
    {"id": "08", "title": "特别说明"},
]

CHAPTER_IDS = {c["id"] for c in REPORT_CHAPTERS}
CHAPTER_BY_ID = {c["id"]: c for c in REPORT_CHAPTERS}

# 品牌色（参考民生银行模板）
COLOR_GOLD = RGBColor(0xC4, 0xA0, 0x52)
COLOR_GOLD_DARK = RGBColor(0x8B, 0x6F, 0x3E)
COLOR_DARK_BG = RGBColor(0x1A, 0x1A, 0x1C)
COLOR_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_FOOTER = RGBColor(0x33, 0x33, 0x33)
COLOR_BADGE = RGBColor(0x8B, 0x6F, 0x3E)


def _blank_layout(prs: Presentation):
    return prs.slide_layouts[6]


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return box


def _add_cover_decor(slide, width, height) -> None:
    """封面右侧金色几何装饰（简化）。"""
    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        width - Inches(4.2),
        height - Inches(3.6),
        Inches(4.5),
        Inches(3.2),
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = COLOR_GOLD
    tri.line.fill.background()
    tri.rotation = 180.0

    tri2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        width - Inches(2.8),
        Inches(0.2),
        Inches(2.4),
        Inches(1.8),
    )
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = COLOR_GOLD_DARK
    tri2.line.fill.background()
    tri2.rotation = 180.0


def _add_cover_slide(
    prs: Presentation,
    *,
    customer_name: str,
    branch_name: str,
    advisor_name: str,
    contact_phone: str,
    report_date: date,
) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_slide_bg(slide, COLOR_DARK_BG)
    w, h = prs.slide_width, prs.slide_height
    _add_cover_decor(slide, w, h)

    _add_textbox(
        slide, Inches(0.6), Inches(0.45), Inches(5), Inches(0.5),
        "中国民生银行  CHINA MINSHENG BANK",
        font_size=11, color=COLOR_GOLD,
    )
    _add_textbox(
        slide, Inches(0.6), Inches(2.0), Inches(8), Inches(0.9),
        f"致：{customer_name}",
        font_size=36, bold=True, color=COLOR_GOLD,
    )
    _add_textbox(
        slide, Inches(0.6), Inches(2.95), Inches(9), Inches(1.2),
        "资产分析报告",
        font_size=52, bold=True, color=COLOR_GOLD,
    )

    date_str = report_date.strftime("%Y.%m.%d")
    meta_lines = [
        branch_name or "--",
        advisor_name or "--",
        f"联系电话：{contact_phone or '--'}",
        f"报告日期：{date_str}",
    ]
    y = h - Inches(2.3)
    for line in meta_lines:
        _add_textbox(
            slide, Inches(0.6), y, Inches(6), Inches(0.35),
            line, font_size=12, color=COLOR_GOLD,
        )
        y += Inches(0.38)


def _add_contents_header(slide, width) -> None:
    _add_textbox(
        slide, Inches(0.55), Inches(0.45), Inches(4), Inches(0.7),
        "目录/Contents",
        font_size=28, bold=True, color=COLOR_TEXT,
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55), Inches(1.15), width - Inches(1.1), Pt(1.5),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()

    _add_textbox(
        slide, width - Inches(4.5), Inches(0.42), Inches(4), Inches(0.35),
        "中国民生银行",
        font_size=10, color=COLOR_MUTED, align=PP_ALIGN.RIGHT,
    )
    _add_textbox(
        slide, width - Inches(4.5), Inches(0.72), Inches(4), Inches(0.35),
        "资产分析报告",
        font_size=14, bold=True, color=COLOR_GOLD_DARK, align=PP_ALIGN.RIGHT,
    )


def _add_chapter_badge(slide, left, top, chapter_id: str, title: str) -> None:
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.42), Inches(0.42),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_BADGE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = chapter_id
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Microsoft YaHei"

    _add_textbox(
        slide, left + Inches(0.52), top + Inches(0.04), Inches(2.2), Inches(0.38),
        title, font_size=16, bold=True, color=COLOR_TEXT,
    )


def _add_contents_slide(prs: Presentation, selected_ids: list[str]) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_slide_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
    w, h = prs.slide_width, prs.slide_height
    _add_contents_header(slide, w)

    cols = 3
    col_w = Inches(3.8)
    start_x = Inches(0.55)
    start_y = Inches(1.55)
    row_h = Inches(0.75)

    for idx, cid in enumerate(selected_ids):
        ch = CHAPTER_BY_ID[cid]
        col = idx % cols
        row = idx // cols
        x = start_x + col * col_w
        y = start_y + row * row_h
        _add_chapter_badge(slide, x, y, ch["id"], ch["title"])

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), h - Inches(0.35), w, Inches(0.35),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = COLOR_FOOTER
    footer.line.fill.background()


def _add_blank_chapter_slide(prs: Presentation, chapter_id: str, title: str) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_slide_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
    _add_textbox(
        slide, Inches(0.55), Inches(0.45), Inches(8), Inches(0.5),
        f"{chapter_id}  {title}",
        font_size=20, bold=True, color=COLOR_TEXT,
    )


def normalize_chapters(selected: Iterable[str]) -> list[str]:
    ids = []
    seen = set()
    for cid in selected:
        if cid not in CHAPTER_IDS or cid in seen:
            continue
        ids.append(cid)
        seen.add(cid)
    return [c["id"] for c in REPORT_CHAPTERS if c["id"] in seen]


def build_allocation_report_ppt(
    *,
    customer_name: str,
    selected_chapters: list[str],
    report_date: date | None = None,
    branch_name: str = "--",
    advisor_name: str = "--",
    contact_phone: str = "--",
) -> bytes:
    chapters = normalize_chapters(selected_chapters)
    if not chapters:
        raise ValueError("请至少选择一个章节")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rd = report_date or date.today()
    _add_cover_slide(
        prs,
        customer_name=customer_name,
        branch_name=branch_name,
        advisor_name=advisor_name,
        contact_phone=contact_phone,
        report_date=rd,
    )
    _add_contents_slide(prs, chapters)
    for cid in chapters:
        ch = CHAPTER_BY_ID[cid]
        _add_blank_chapter_slide(prs, ch["id"], ch["title"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def report_filename(customer_name: str) -> str:
    safe = "".join(c for c in customer_name if c not in '\\/:*?"<>|').strip() or "客户"
    return f"资产配置报告-{safe}.pptx"


def content_disposition_attachment(filename: str) -> str:
    """HTTP 头仅支持 latin-1，中文文件名用 RFC 5987 filename*。"""
    from urllib.parse import quote

    ascii_name = "allocation-report.pptx"
    encoded = quote(filename, safe="")
    return f'attachment; filename="{ascii_name}"; filename*=UTF-8\'\'{encoded}'
