"""Tests for allocation report PPT export."""

import io

import pytest
from pptx import Presentation

from core.allocation_report_ppt import (
    CHAPTER_IDS,
    build_allocation_report_ppt,
    normalize_chapters,
    report_filename,
)


class TestAllocationReportPpt:
    def test_normalize_chapters_preserves_order(self):
        assert normalize_chapters(["03", "01", "99"]) == ["01", "03"]

    def test_build_ppt_slide_count(self):
        data = build_allocation_report_ppt(
            customer_name="赵女士",
            selected_chapters=["01", "02", "03"],
        )
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 5  # cover + contents + 3 chapters

    def test_build_ppt_all_chapters(self):
        data = build_allocation_report_ppt(
            customer_name="张女士",
            selected_chapters=list(CHAPTER_IDS),
        )
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 2 + len(CHAPTER_IDS)

    def test_empty_chapters_raises(self):
        with pytest.raises(ValueError):
            build_allocation_report_ppt(customer_name="测试", selected_chapters=[])

    def test_report_filename(self):
        assert report_filename("赵女士") == "资产配置报告-赵女士.pptx"

    def test_content_disposition_ascii_safe(self):
        from core.allocation_report_ppt import content_disposition_attachment

        header = content_disposition_attachment("资产配置报告-赵女士.pptx")
        assert "filename=\"allocation-report.pptx\"" in header
        assert "filename*=UTF-8" in header
        header.encode("latin-1")
